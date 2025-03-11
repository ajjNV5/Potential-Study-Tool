from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Define brand colors
brand_colors = {
    "primary": RGBColor(0, 98, 152),
    "secondary_blue": RGBColor(113, 178, 201),
    "secondary_gray": RGBColor(127, 127, 127),
    "tertiary_blue": RGBColor(1, 66, 106),
    "tertiary_green": RGBColor(116, 170, 80),
    "tertiary_orange": RGBColor(115, 111, 33),
    "tertiary_gray": RGBColor(183, 169, 154)
}

# Define brand fonts
brand_fonts = {
    "headings_bold": "Franklin Gothic Demi",
    "headings_std": "Franklin Gothic Medium",
    "body": "Franklin Gothic Book",
    "graphics": "Roboto"
}

# Create a presentation object
prs = Presentation()

# Title Slide Layout
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Title Slide"
title.text_frame.paragraphs[0].font.name = brand_fonts["headings_bold"]
title.text_frame.paragraphs[0].font.size = Pt(44)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

subtitle.text = "Subtitle"
subtitle.text_frame.paragraphs[0].font.name = brand_fonts["headings_std"]
subtitle.text_frame.paragraphs[0].font.size = Pt(32)
subtitle.text_frame.paragraphs[0].font.color.rgb = brand_colors["secondary_gray"]

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = brand_colors["primary"]

# Content Slide Layout
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Content Slide"
title.text_frame.paragraphs[0].font.name = brand_fonts["headings_bold"]
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.color.rgb = brand_colors["tertiary_blue"]

content.text = "Content goes here..."
content.text_frame.paragraphs[0].font.name = brand_fonts["body"]
content.text_frame.paragraphs[0].font.size = Pt(24)
content.text_frame.paragraphs[0].font.color.rgb = brand_colors["secondary_gray"]

# Section Header Slide Layout
slide_layout = prs.slide_layouts[2]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title

title.text = "Section Header"
title.text_frame.paragraphs[0].font.name = brand_fonts["headings_bold"]
title.text_frame.paragraphs[0].font.size = Pt(44)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = brand_colors["tertiary_blue"]

# Data Slide Layout
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title

title.text = "Data Slide"
title.text_frame.paragraphs[0].font.name = brand_fonts["headings_bold"]
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.color.rgb = brand_colors["primary"]

# Closing Slide Layout
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Closing Slide"
title.text_frame.paragraphs[0].font.name = brand_fonts["headings_bold"]
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.color.rgb = brand_colors["primary"]

content.text = "Thank you!"
content.text_frame.paragraphs[0].font.name = brand_fonts["body"]
content.text_frame.paragraphs[0].font.size = Pt(24)
content.text_frame.paragraphs[0].font.color.rgb = brand_colors["secondary_gray"]

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = brand_colors["tertiary_gray"]

# Save the presentation as a .pptx file
prs.save("data_team_template.pptx")

print("PowerPoint template 'data_team_template.pptx' has been created successfully.")